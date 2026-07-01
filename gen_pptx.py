"""Generate presentation.pptx from the project's slide content."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BG   = RGBColor(0x0D, 0x1B, 0x2A)   # near-black navy
ACCENT    = RGBColor(0xE0, 0x3A, 0x3A)   # red
BLUE      = RGBColor(0x4A, 0x9E, 0xFF)   # blue
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GREY      = RGBColor(0xB0, 0xB8, 0xC8)
GREEN     = RGBColor(0x3D, 0xD6, 0x8C)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

TAGLE_IMG = "tagle_tag..png"


def prs() -> Presentation:
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p


def blank(p: Presentation):
    return p.slides.add_slide(p.slide_layouts[6])   # completely blank


def bg(slide, color: RGBColor = DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txb(slide, text, l, t, w, h, size=24, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def line(slide, x1, y1, x2, y2, color=ACCENT, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,   # straight
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width     = width


def bullet_box(slide, items, l, t, w, h, size=20, color=WHITE, icon="•"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.space_before = Pt(4)
        run = para.add_run()
        run.text = f"{icon}  {item}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color


# ── SLIDE 1 — Title ──────────────────────────────────────────────────────────
def slide_title(p: Presentation):
    s = blank(p)
    bg(s)

    # red accent bar on left
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), SLIDE_H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # Title
    txb(s, "Intelligent Observability\n& Event Watchdog",
        0.3, 0.5, 8.5, 2.2, size=40, bold=True, color=WHITE)

    # Sub-title
    txb(s, "A Python / FastAPI AI-Directed SRE Tool",
        0.3, 2.9, 8.5, 0.7, size=22, color=GREY, italic=True)

    # Author + Tagle text
    txb(s, "Aditya Patil",
        0.3, 3.8, 6, 0.5, size=20, bold=True, color=WHITE)
    txb(s, "The Connector  ·  AI Readiness Type  ·  Tagle.ai",
        0.3, 4.35, 8, 0.5, size=16, color=BLUE)
    txb(s, '"You bring people together to make AI work for everyone"',
        0.3, 4.85, 8.5, 0.5, size=14, italic=True, color=GREY)

    # Bottom tagline
    txb(s, "Zero cloud dependencies  ·  SQLite  ·  FastAPI  ·  Docker-ready  ·  Built < 3 hours",
        0.3, 6.6, 10, 0.5, size=13, color=GREY)

    # Tagle tag image (right side)
    if os.path.exists(TAGLE_IMG):
        try:
            s.shapes.add_picture(TAGLE_IMG, Inches(9.5), Inches(1.5),
                                 Inches(3.5), Inches(4.5))
        except Exception:
            pass   # image exists but can't be embedded — skip silently


# ── SLIDE 2 — Problem ─────────────────────────────────────────────────────────
def slide_problem(p: Presentation):
    s = blank(p)
    bg(s)
    txb(s, "The Problem", 0.6, 0.3, 12, 0.8, size=34, bold=True, color=ACCENT)

    txb(s, "Modern systems generate more logs than humans can monitor",
        0.6, 1.25, 12, 0.6, size=20, color=GREY, italic=True)

    sources = [
        "Apache HTTP Server   →   error.log   →   thousands of lines / hour",
        "BGL Supercomputer    →   131,072 processors   →   millions of kernel events",
        "Linux syslog         →   SSH attacks, FTP floods   →   constant noise",
    ]
    bullet_box(s, sources, 0.6, 2.0, 12, 1.5, size=17, color=BLUE, icon="▶")

    challenges = [
        "Alert fatigue — ops teams drown in undifferentiated noise",
        "Invisible spikes — real error bursts buried in background chatter",
        "Delayed response — anomalies spotted too late, after impact",
        "Fragmented tooling — parse here, detect there, alert somewhere else",
        "Cost — Datadog / Splunk / New Relic charge thousands per month",
    ]
    bullet_box(s, challenges, 0.6, 3.6, 12, 2.5, size=17, icon="✗", color=WHITE)

    txb(s, "→  One self-contained Python service: ingest · detect · alert · visualise",
        0.6, 6.6, 12, 0.5, size=15, bold=True, color=GREEN)


# ── SLIDE 3 — Architecture ────────────────────────────────────────────────────
def slide_arch(p: Presentation):
    s = blank(p)
    bg(s)
    txb(s, "Architecture", 0.6, 0.3, 12, 0.8, size=34, bold=True, color=ACCENT)

    mono = (
        "┌───────────────────── Single Process (FastAPI + Gunicorn) ─────────────────────┐\n"
        "│  9 REST Endpoints           WebSocket /api/stream                              │\n"
        "│  /api/logs                  Log Parser  (log_parser.py)                        │\n"
        "│  /api/analyze               apache_logs.csv  →  unified EventCreate            │\n"
        "│  /api/anomalies             bgl_logs.csv     →  unified EventCreate            │\n"
        "│  /api/webhooks              linux_logs.csv   →  unified EventCreate            │\n"
        "│                                    ↓                                           │\n"
        "│                     SQLite  watchdog.db  (WAL mode)                            │\n"
        "│                     logs  ·  anomalies  ·  webhook_alerts                      │\n"
        "│                                    ↓                                           │\n"
        "│              Anomaly Detector  (Z-Score  ·  IQR  ·  Moving Average)            │\n"
        "│                                    ↓                                           │\n"
        "│              Webhook Handler  →  WH-<uuid> payload  →  SQLite                  │\n"
        "│                                                                                 │\n"
        "│   Dashboard  static/index.html  ·  Chart.js CDN  ·  WebSocket JS              │\n"
        "└─────────────────────────────────────────────────────────────────────────────────┘"
    )
    box = s.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(12.5), Inches(5.5))
    tf  = box.text_frame
    tf.word_wrap = False
    p_  = tf.paragraphs[0]
    run = p_.add_run()
    run.text = mono
    run.font.size      = Pt(11)
    run.font.color.rgb = BLUE
    run.font.name      = "Courier New"

    txb(s, "One command:  uvicorn main:app --reload --port 8000",
        0.6, 6.9, 12, 0.4, size=14, color=GREEN, bold=True)


# ── SLIDE 4 — Key Features ────────────────────────────────────────────────────
def slide_features(p: Presentation):
    s = blank(p)
    bg(s)
    txb(s, "Key Features", 0.6, 0.3, 12, 0.8, size=34, bold=True, color=ACCENT)

    features = [
        ("1. Multi-Source Log Ingest",
         "Apache · BGL Supercomputer · Linux syslog → unified {timestamp, severity, source, message} schema"),
        ("2. Three Anomaly Detection Algorithms",
         "Z-Score (primary)  ·  IQR  ·  Moving Average — configurable window & threshold"),
        ("3. Structured Webhook Alerting",
         "WH-<uuid> JSON payload per anomaly: CRITICAL (z>3) PAGE_ON_CALL · WARNING NOTIFY"),
        ("4. Real-time Dashboard",
         "Health pill · 4 metric cards · Error Rate chart · Severity bar · Source doughnut · Live stream"),
        ("5. Production Docker Setup",
         "Multi-stage image · non-root user · named volume · dev hot-reload override"),
    ]
    y = 1.4
    for title, body in features:
        txb(s, title, 0.6, y, 12, 0.45, size=18, bold=True, color=WHITE)
        txb(s, body,  0.9, y + 0.45, 11.5, 0.5, size=15, color=GREY)
        y += 1.12


# ── SLIDE 5 — Live Demo Flow ──────────────────────────────────────────────────
def slide_demo(p: Presentation):
    s = blank(p)
    bg(s)
    txb(s, "Live Demo Flow", 0.6, 0.3, 12, 0.8, size=34, bold=True, color=ACCENT)

    steps = [
        ("Step 1", "uvicorn main:app --reload --port 8000 → auto-ingest 160 logs → open localhost:8000"),
        ("Step 2", "Observe pre-loaded anomalies: Apache spike · BGL FATAL burst · Linux SSH storm"),
        ("Step 3", "Switch method to IQR → Run Analysis  |  Lower Z threshold → more alerts surface"),
        ("Step 4", "Click Simulate Live Feed → watch chart grow at 300 ms/log → Stop"),
        ("Step 5", "python synthetic_data_gen.py --rows 500  →  Ingest  →  Run Analysis"),
        ("Step 6", "Expand a SPIKE badge → WH-<uuid> JSON with action: PAGE_ON_CALL"),
        ("Step 7", "Visit /docs → full OpenAPI UI  |  /api/webhooks → raw JSON history"),
    ]
    y = 1.4
    for label, desc in steps:
        txb(s, label, 0.6, y, 1.8, 0.45, size=16, bold=True, color=ACCENT)
        txb(s, desc,  2.5, y, 10.3, 0.45, size=16, color=WHITE)
        y += 0.82


# ── SLIDE 6 — Tech Stack & AI-Directed ───────────────────────────────────────
def slide_tech(p: Presentation):
    s = blank(p)
    bg(s)
    txb(s, "Tech Stack & AI-Directed Development", 0.6, 0.3, 12, 0.8,
        size=30, bold=True, color=ACCENT)

    stack = [
        "API:        FastAPI 0.115  —  async, auto-docs, WebSocket native",
        "Database:   SQLite + WAL   —  zero-config, zero cost, ACID compliant",
        "Server:     Gunicorn + uvicorn worker  —  production lifecycle",
        "Frontend:   Vanilla JS + Chart.js CDN  —  no build step, single HTML",
        "Detection:  Pure Python stdlib math     —  no NumPy dependency",
        "Containers: Docker + Compose            —  reproducible deploys",
    ]
    bullet_box(s, stack, 0.6, 1.2, 8.5, 2.8, size=16, icon="▶", color=BLUE)

    txb(s, "How the session worked:", 0.6, 4.1, 8, 0.5, size=18, bold=True, color=WHITE)
    loop = [
        "Human defines requirements  →  AI translates to architecture",
        "Human specifies CSV schemas →  AI writes complete implementation",
        "Human reviews screenshot    →  AI diagnoses bug, rewrites files",
        "Human approves each phase   →  AI executes with zero manual edits",
    ]
    bullet_box(s, loop, 0.6, 4.65, 8.5, 2.1, size=15, icon="→", color=GREY)

    txb(s, "Full MVP · 22 files · ~1 800 lines · < 3 hours · 0 cloud services · 0 manual edits",
        0.6, 6.85, 12, 0.4, size=13, bold=True, color=GREEN)


# ── SLIDE 7 — Future Improvements ────────────────────────────────────────────
def slide_future(p: Presentation):
    s = blank(p)
    bg(s)
    txb(s, "Future Improvements", 0.6, 0.3, 12, 0.8, size=34, bold=True, color=ACCENT)

    near = [
        "ML anomaly detection — Isolation Forest / LSTM for non-stationary series",
        "Real webhook delivery — POST to Slack / PagerDuty with retry queue",
        "Alert deduplication — suppress repeat alerts for same window",
        "Rule engine — user-defined threshold rules (error_rate > 60% for 10 min)",
        "JWT / API-key authentication on all endpoints",
    ]
    txb(s, "Near-term", 0.6, 1.25, 6, 0.45, size=18, bold=True, color=WHITE)
    bullet_box(s, near, 0.6, 1.75, 6.2, 2.5, size=15, icon="•", color=GREY)

    medium = [
        "PostgreSQL — multi-worker horizontal scaling",
        "Log tail mode — watch live files via inotify",
        "Grafana integration — Prometheus /metrics endpoint",
        "Email digest — daily anomaly summaries via SMTP",
    ]
    txb(s, "Medium-term", 6.9, 1.25, 6, 0.45, size=18, bold=True, color=WHITE)
    bullet_box(s, medium, 6.9, 1.75, 6, 2.0, size=15, icon="•", color=GREY)

    txb(s, "Vision: Any log format  ·  LLM root cause analysis  ·  K8s + HPA  ·  TimescaleDB",
        0.6, 6.5, 12, 0.6, size=15, italic=True, color=BLUE)


# ── SLIDE 8 — Summary ─────────────────────────────────────────────────────────
def slide_summary(p: Presentation):
    s = blank(p)
    bg(s)
    txb(s, "Summary", 0.6, 0.3, 12, 0.8, size=34, bold=True, color=ACCENT)

    rows = [
        ("Lines of code",        "~1 800 (Python + HTML + YAML)"),
        ("Files delivered",      "22"),
        ("Build time",           "< 3 hours  (AI-directed session)"),
        ("External paid services","0"),
        ("Cloud dependencies",   "0"),
        ("Manual edits by human","0"),
        ("Bugs fixed",           "1  (empty dashboard → auto-ingest on startup)"),
        ("Anomaly methods",      "3  (Z-Score · IQR · Moving Average)"),
        ("Docker image size",    "~200 MB  (multi-stage slim)"),
    ]

    y = 1.35
    for label, value in rows:
        # alternating row background
        box = s.shapes.add_shape(1, Inches(0.5), Inches(y - 0.05),
                                 Inches(12.3), Inches(0.52))
        box.fill.solid()
        box.fill.fore_color.rgb = (RGBColor(0x14, 0x2A, 0x40)
                                   if rows.index((label, value)) % 2 == 0
                                   else DARK_BG)
        box.line.fill.background()

        txb(s, label, 0.6, y, 6.5, 0.45, size=16, color=GREY)
        txb(s, value, 7.2, y, 5.5, 0.45, size=16, bold=True, color=WHITE)
        y += 0.55

    txb(s, "Built by Aditya Patil  ·  The Connector (Tagle.ai)  ·  July 2026",
        0.6, 7.05, 12, 0.35, size=13, color=GREY,
        align=PP_ALIGN.CENTER)


# ── MAIN ──────────────────────────────────────────────────────────────────────
def main():
    p = prs()
    slide_title(p)
    slide_problem(p)
    slide_arch(p)
    slide_features(p)
    slide_demo(p)
    slide_tech(p)
    slide_future(p)
    slide_summary(p)
    p.save("presentation.pptx")
    print(f"Saved presentation.pptx  ({len(p.slides)} slides)")


if __name__ == "__main__":
    main()
